from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x25, 0x63, 0xeb)
LIGHT_BG = RGBColor(0xf8, 0xf9, 0xfa)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x16, 0xa3, 0x4a)
ORANGE = RGBColor(0xea, 0x58, 0x0c)
RED = RGBColor(0xdc, 0x26, 0x26)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox


# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_shape_fill(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), ACCENT)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Veil Stack", font_size=48, bold=True, color=WHITE)
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
             "Decentralized Container Orchestration on FEVM + libp2p",
             font_size=22, color=RGBColor(0xa0, 0xa0, 0xc0))
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
             "NLnet Grant Application  |  MIT License  |  Open Source",
             font_size=16, color=GRAY)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Sumanjeet  &  Shivam Kumar (shivv23)",
             font_size=14, color=RGBColor(0x88, 0x88, 0xaa))

# ── Slide 2: Problem ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             "The Problem", font_size=32, bold=True, color=WHITE)

problems = [
    "Single points of failure \u2014 centralized orchestrators (Kubernetes) can bring down entire stacks",
    "Vendor lock-in \u2014 cloud providers control scheduling, pricing, and data placement",
    "No verifiable storage \u2014 container images sit on opaque registries with no audit trail",
    "No privacy for multi-org clusters \u2014 organizations must expose resource metrics to a shared control plane",
]
add_bullet_list(slide, Inches(1), Inches(1.8), Inches(11), Inches(4.5),
                [f"\u2022  {p}" for p in problems], font_size=18, color=DARK)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
             "Container orchestration needs a decentralized, verifiable, and privacy-preserving alternative.",
             font_size=16, bold=True, color=ACCENT)

# ── Slide 3: Solution ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             "Our Solution", font_size=32, bold=True, color=WHITE)

solutions = [
    ("FEVM Smart Contract", "On-chain governance: member management, image registry, replica balancing, deal anchoring"),
    ("libp2p Networking", "Decentralized cluster: TCP transport, Noise encryption, GossipSub heartbeats \u2014 no central control plane"),
    ("Docker Runtime", "Pull, create, start, stop containers via Docker Engine API, driven by on-chain events"),
    ("IPFS + Filecoin", "Deployment manifests pinned to IPFS; V2 adds automated Filecoin storage deals per workload"),
    ("FHE Scheduling (Planned)", "Encrypted telemetry via Zama FHE SDK for regulated workloads"),
]

for i, (title, desc) in enumerate(solutions):
    y = 1.6 + i * 1.1
    add_text_box(slide, Inches(1), Inches(y), Inches(3.5), Inches(0.5),
                 title, font_size=16, bold=True, color=ACCENT)
    add_text_box(slide, Inches(4.5), Inches(y), Inches(8), Inches(0.5),
                 desc, font_size=14, color=DARK)

# ── Slide 4: Current State (V1) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             "Current State (V1 \u2014 Working)", font_size=32, bold=True, color=WHITE)

v1_items = [
    ("Canteen.sol on FEVM Calibration", "Deployed & verified", GREEN),
    ("Web Dashboard (React + D3)", "Live on Vercel", GREEN),
    ("MetaMask Integration", "Filecoin Calibration (chain 314159)", GREEN),
    ("libp2p Cluster Networking", "TCP, Noise, mDNS, GossipSub", GREEN),
    ("Docker Container Runtime", "Pull, create, start, stop", GREEN),
    ("Event-Driven Scheduler", "MemberJoin, MemberLeave, MemberImageUpdate", GREEN),
    ("IPFS Pinning (Pinata)", "Deployment manifests", GREEN),
    ("CI/CD Pipeline", "GitHub Actions: contract tests + Docker build", GREEN),
]

for i, (name, status, color) in enumerate(v1_items):
    y = 1.5 + i * 0.7
    add_text_box(slide, Inches(1), Inches(y), Inches(0.3), Inches(0.4),
                 "\u2713", font_size=16, bold=True, color=color)
    add_text_box(slide, Inches(1.4), Inches(y), Inches(4.5), Inches(0.4),
                 name, font_size=15, bold=True, color=DARK)
    add_text_box(slide, Inches(6), Inches(y), Inches(6), Inches(0.4),
                 status, font_size=14, color=GRAY)

# ── Slide 5: Architecture ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             "Architecture", font_size=32, bold=True, color=WHITE)

layers = [
    ("Dashboard (React + D3 + Web3)", 5.5, 1.8, ACCENT),
    ("FEVM Contract (Canteen.sol)", 5.5, 3.0, RGBColor(0x7c, 0x3a, 0xed)),
    ("Veil Nodes (libp2p) + Scheduler", 5.5, 4.2, GREEN),
    ("Docker Hosts", 5.5, 5.4, ORANGE),
    ("Filecoin Network (Calibration \u2192 Mainnet)", 5.5, 6.4, RGBColor(0x0e, 0x7a, 0x0e)),
]

for label, x, y, color in layers:
    w = Inches(6)
    h = Inches(0.7)
    shape = add_shape_fill(slide, Inches(x - 3), Inches(y), w, h, color)
    shape.text_frame.paragraphs[0].text = label
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Calibri"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Arrows between layers
for y_pos in [2.55, 3.75, 4.95, 6.1]:
    add_text_box(slide, Inches(5.3), Inches(y_pos), Inches(0.5), Inches(0.4),
                 "\u25BC", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Side notes
side_notes = [
    (1.2, 1.8, "React + D3\nforce graph\nWeb3/MetaMask"),
    (1.2, 3.0, "Membership\nImage registry\nReplica balancing"),
    (1.2, 4.2, "TCP + Noise\nGossipSub\nmDNS discovery"),
    (1.2, 5.4, "Docker Engine\nAPI\nSocket proxy"),
    (1.2, 6.4, "Storage deals\nCID verification\nMainnet migration"),
]
for x, y, note in side_notes:
    add_text_box(slide, Inches(x), Inches(y), Inches(3.5), Inches(0.7),
                 note, font_size=10, color=GRAY)

# ── Slide 6: Filecoin Integration ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             "Filecoin Integration Roadmap", font_size=32, bold=True, color=WHITE)

# V1 column
add_text_box(slide, Inches(1), Inches(1.5), Inches(5.5), Inches(0.5),
             "V1 (Deployed)", font_size=20, bold=True, color=GREEN)
v1_features = [
    "Canteen.sol on FEVM Calibration",
    "Member management + image registry",
    "Event-driven scheduler",
    "IPFS pinning via Pinata",
    "Dashboard with MetaMask",
]
add_bullet_list(slide, Inches(1), Inches(2.1), Inches(5.5), Inches(3),
                [f"\u2713  {f}" for f in v1_features], font_size=14, color=DARK)

# V2 column
add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
             "V2 (Planned \u2014 This Grant)", font_size=20, bold=True, color=ACCENT)
v2_features = [
    "StorageDeal struct on-chain",
    "Lotus JSON-RPC deal proposal",
    "Deal lifecycle monitoring",
    "CID-verified image retrieval",
    "Multi-provider fallback",
]
add_bullet_list(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(3),
                [f"\u25CB  {f}" for f in v2_features], font_size=14, color=DARK)

# V3 column
add_text_box(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(0.5),
             "V3 (Research)", font_size=20, bold=True, color=ORANGE)
v3_features = [
    "FHE confidential scheduling (Zama)",
    "10-node CI + federation model",
    "Mainnet migration",
    "Security audit",
]
add_bullet_list(slide, Inches(7), Inches(5.1), Inches(5.5), Inches(2.5),
                [f"\u25CB  {f}" for f in v3_features], font_size=14, color=DARK)

# ── Slide 7: Milestones ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             "Funding-Led Milestones", font_size=32, bold=True, color=WHITE)

milestones = [
    ("M1: Filecoin Deal Pipeline", "\u20AC15,000", "Months 1-5",
     "StorageDeal struct, Lotus JSON-RPC, deal lifecycle, dashboard deal tab"),
    ("M2: CID Verification + Multi-Provider", "\u20AC12,000", "Months 1-4",
     "CID-verified retrieval, provider fallback, retry logic, E2E tests"),
    ("M3: FHE Confidential Scheduling", "\u20AC10,000", "Months 1-6",
     "Zama FHE SDK, encrypted telemetry, ciphertext scheduling, benchmarks"),
    ("M4: Production Hardening", "\u20AC8,000", "Months 1-5",
     "Security audit, 10-node CI, federation model, mainnet migration plan"),
]

for i, (title, amount, timeline, desc) in enumerate(milestones):
    y = 1.5 + i * 1.4
    add_shape_fill(slide, Inches(0.8), Inches(y), Inches(11.7), Inches(1.2), LIGHT_BG)
    add_text_box(slide, Inches(1), Inches(y + 0.05), Inches(5), Inches(0.4),
                 title, font_size=16, bold=True, color=DARK)
    add_text_box(slide, Inches(6.5), Inches(y + 0.05), Inches(2), Inches(0.4),
                 amount, font_size=16, bold=True, color=ACCENT)
    add_text_box(slide, Inches(9), Inches(y + 0.05), Inches(3), Inches(0.4),
                 timeline, font_size=14, color=GRAY)
    add_text_box(slide, Inches(1), Inches(y + 0.5), Inches(11), Inches(0.6),
                 desc, font_size=13, color=GRAY)

# Total
add_shape_fill(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4), ACCENT)
add_text_box(slide, Inches(1), Inches(7.0), Inches(5), Inches(0.4),
             "Total", font_size=14, bold=True, color=WHITE)
add_text_box(slide, Inches(6.5), Inches(7.0), Inches(2), Inches(0.4),
             "\u20AC45,000", font_size=14, bold=True, color=WHITE)

# ── Slide 8: Use Cases ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             "Use Cases", font_size=32, bold=True, color=WHITE)

use_cases = [
    ("Decentralized Cloud Compute",
     "On-chain governed container orchestration with libp2p cluster networking. No single point of failure."),
    ("Regulated Workloads (Healthcare, Defense)",
     "FHE-encrypted scheduling metrics. Audit trail on-chain via FEVM. Zero-trust multi-org clusters."),
    ("Cross-Org Compute Cooperatives",
     "libp2p federation + FEVM governance enable multi-org clusters without trust assumptions."),
    ("AI/ML Training Pipelines",
     "Large model artifacts stored on Filecoin with CID verification. Programmatic deal origination."),
]

for i, (title, desc) in enumerate(use_cases):
    y = 1.6 + i * 1.3
    add_text_box(slide, Inches(1), Inches(y), Inches(11), Inches(0.4),
                 title, font_size=17, bold=True, color=ACCENT)
    add_text_box(slide, Inches(1), Inches(y + 0.45), Inches(11), Inches(0.6),
                 desc, font_size=14, color=GRAY)

# ── Slide 9: Links & Close ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_shape_fill(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), ACCENT)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Thank You", font_size=42, bold=True, color=WHITE)

links = [
    "Repository:  github.com/shivv23/Veil-Stack",
    "Contract:  0x04dEf60e2853E4d654b366cd8103F929c456d4b7 (FEVM Calibration)",
    "Dashboard:  veil-stack-canteen.vercel.app/dashboard/",
    "License:  MIT",
    "Contact:  shivv23 (GitHub)",
]
add_bullet_list(slide, Inches(1), Inches(3.6), Inches(11), Inches(3),
                links, font_size=16, color=RGBColor(0xcc, 0xcc, 0xdd))

prs.save(r"C:\Users\HP\Desktop\VeilSTack\Veil-Stack\Veil_Stack_NLnet.pptx")
print("PPT saved.")
